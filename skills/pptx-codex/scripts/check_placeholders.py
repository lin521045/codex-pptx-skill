#!/usr/bin/env python3
"""Scan a PPTX deck for common placeholder text."""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation


DEFAULT_PATTERNS = [
    r"\blorem\b",
    r"\bipsum\b",
    r"\btodo\b",
    r"\btbd\b",
    r"\bxxx+\b",
    r"click to add",
    r"your (title|name|text)",
    r"placeholder",
]


def main() -> int:
    parser = argparse.ArgumentParser(description="Detect placeholder text in a PPTX deck.")
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument(
        "--pattern",
        action="append",
        default=[],
        help="Extra regex pattern to scan for. Can be used more than once.",
    )
    args = parser.parse_args()

    presentation = Presentation(args.pptx)
    patterns = [re.compile(pattern, re.IGNORECASE) for pattern in DEFAULT_PATTERNS + args.pattern]
    found = False

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            for pattern in patterns:
                for match in pattern.finditer(text):
                    found = True
                    print(
                        f"Slide {slide_index}: matched '{match.group(0)}' "
                        f"with pattern '{pattern.pattern}'"
                    )

    if found:
        return 1

    print("No placeholder text detected.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

