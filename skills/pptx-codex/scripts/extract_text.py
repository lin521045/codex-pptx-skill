#!/usr/bin/env python3
"""Extract readable text from a PPTX deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def paragraphs_from_shape(shape) -> list[str]:
    lines: list[str] = []

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs) or paragraph.text
            text = text.strip()
            if not text:
                continue
            prefix = ("  " * paragraph.level + "- ") if paragraph.level >= 0 else ""
            lines.append(f"{prefix}{text}")

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            values = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
            if any(values):
                lines.append("| " + " | ".join(values) + " |")

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            lines.extend(paragraphs_from_shape(child))

    return lines


def slide_title(slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is None:
        return "Untitled slide"
    title = title_shape.text.strip()
    return title or "Untitled slide"


def notes_text(slide) -> list[str]:
    try:
        frame = slide.notes_slide.notes_text_frame
    except Exception:
        return []
    content = [line.strip() for line in frame.text.splitlines() if line.strip()]
    return content


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract text from a PowerPoint presentation.")
    parser.add_argument("pptx", type=Path, help="Path to the .pptx file")
    parser.add_argument("--include-notes", action="store_true", help="Include speaker notes")
    args = parser.parse_args()

    presentation = Presentation(args.pptx)

    print(f"# {args.pptx.name}")
    print()
    for index, slide in enumerate(presentation.slides, start=1):
        print(f"## Slide {index}: {slide_title(slide)}")
        body: list[str] = []
        for shape in slide.shapes:
            body.extend(paragraphs_from_shape(shape))
        if body:
            for line in body:
                print(line)
        else:
            print("(No extractable text)")
        if args.include_notes:
            notes = notes_text(slide)
            if notes:
                print()
                print("Notes:")
                for line in notes:
                    print(f"- {line}")
        print()

    return 0


if __name__ == "__main__":
    raise SystemExit(main())

